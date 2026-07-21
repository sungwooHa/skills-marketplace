"""
HTML 슬라이드 → PPTX (스크린샷 + 노트 임베드 단일 방식)

원칙
- HTML이 정답. PPTX는 청중 배포/오프라인 보조 산출물.
- 슬라이드 본체: Playwright로 1920×1080 스크린샷을 PNG로 떠서 PPT에 풀 블리드 임베드
- 발표 노트: 같은 디렉토리의 스크립트_화면.html에서 .page 슬라이드의 .script를 추출
            → PPTX 노트 영역(notes_slide)에 텍스트로 삽입
- 받는 사람이 PPT 본문 텍스트를 직접 편집할 수는 없음 (의도된 트레이드오프)

사용법
  python scripts/build_pptx_screenshot.py output/분기_OKR_리뷰/index.html
    → output/분기_OKR_리뷰/index.pptx 생성

  python scripts/build_pptx_screenshot.py path/to/index.html -o out.pptx
  python scripts/build_pptx_screenshot.py index.html --notes 스크립트_화면.html
  python scripts/build_pptx_screenshot.py index.html --keep-screenshots
  python scripts/build_pptx_screenshot.py index.html --settle-ms 500

선행 조건
  pip install playwright python-pptx beautifulsoup4 Pillow
  playwright install chromium
"""
from __future__ import annotations

import argparse
import http.server
import shutil
import socket
import sys
import tempfile
import threading
from pathlib import Path

from bs4 import BeautifulSoup
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu


# 16:9 widescreen, python-pptx 기본 (10" × 5.625" @ 914400 EMU/inch)
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5143500

# 스크린샷용 뷰포트 (HTML이 1920×1080 fixed)
VIEWPORT_W = 1920
VIEWPORT_H = 1080

# 애니메이션 즉시 완료 → 스크린샷이 최종 상태가 되도록 강제 주입
DISABLE_ANIM_CSS = """
*, *::before, *::after {
  transition-duration: 0s !important;
  transition-delay: 0s !important;
  animation-duration: 0s !important;
  animation-delay: 0s !important;
}
.slide.active.play .fade { opacity: 1 !important; transform: none !important; }
"""


def parse_notes(script_html_path: Path) -> list[dict]:
    """스크립트_화면.html에서 .page 슬라이드의 발화 텍스트 추출.

    반환: [{'title': '슬라이드 제목', 'pnum': 'PAGE 5', 'pact': 'Act 1 · S5', 'script': '...'}, ...]
    cover/actbreak는 제외 — .page만 메인 덱과 1:1 sync.
    """
    if not script_html_path.exists():
        return []
    soup = BeautifulSoup(script_html_path.read_text(encoding="utf-8"), "html.parser")
    pages = soup.select(".slide.page")
    notes: list[dict] = []
    for p in pages:
        title_el = p.select_one(".ptitle")
        pnum_el = p.select_one(".pnum")
        script_el = p.select_one(".script")

        title = title_el.get_text(strip=True) if title_el else ""
        pnum = pnum_el.get_text(strip=True) if pnum_el else ""
        pact = p.get("data-act-name", "")

        if script_el:
            for br in script_el.select("br"):
                br.replace_with("\n")
            script = script_el.get_text()
        else:
            script = ""

        notes.append(
            {
                "title": title,
                "pnum": pnum,
                "pact": pact,
                "script": script.strip(),
            }
        )
    return notes


def _start_http_server(directory: Path) -> tuple[http.server.HTTPServer, int]:
    """HTML 디렉토리를 서빙하는 임시 HTTP 서버를 백그라운드에서 시작."""
    # 빈 포트 자동 할당
    sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    sock.bind(("127.0.0.1", 0))
    port = sock.getsockname()[1]
    sock.close()

    handler = lambda *args, **kwargs: http.server.SimpleHTTPRequestHandler(
        *args, directory=str(directory), **kwargs
    )
    server = http.server.HTTPServer(("127.0.0.1", port), handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    return server, port


def screenshot_slides(html_path: Path, out_dir: Path, settle_ms: int = 300) -> list[Path]:
    """Playwright로 #1, #2, ... 해시 네비게이션해서 슬라이드별 PNG 저장."""
    out_dir.mkdir(parents=True, exist_ok=True)

    # 한글 경로 file:// 호환 문제 방지: 임시 HTTP 서버로 서빙
    serve_dir = html_path.parent.absolute()
    server, port = _start_http_server(serve_dir)
    base_url = f"http://127.0.0.1:{port}/{html_path.name}"
    print(f"🌐 임시 HTTP 서버: http://127.0.0.1:{port} (종료: 자동)", file=sys.stderr)

    screenshots: list[Path] = []

    with sync_playwright() as pw:
        browser = pw.chromium.launch()
        context = browser.new_context(
            viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
            device_scale_factor=2,  # Retina 품질
        )
        page = context.new_page()

        # 첫 진입: 슬라이드 수 파악
        page.goto(base_url, wait_until="networkidle")
        # JS가 슬라이드 수를 #total에 채우길 기다림
        page.wait_for_function("document.getElementById('total') && document.getElementById('total').textContent !== '0'")
        total = page.evaluate("document.querySelectorAll('#deck .slide').length")
        if total == 0:
            raise RuntimeError(f"슬라이드를 찾을 수 없음: {html_path}")

        for i in range(1, total + 1):
            page.goto(f"{base_url}#{i}", wait_until="networkidle")
            # 애니메이션 무력화 CSS 주입 (매 페이지 로드 후)
            page.add_style_tag(content=DISABLE_ANIM_CSS)
            # Wait for webfont load (external .otf fetch)
            try:
                page.evaluate("() => document.fonts && document.fonts.ready")
            except Exception:
                pass
            page.wait_for_timeout(settle_ms)
            png_path = out_dir / f"slide-{i:02d}.png"
            page.screenshot(path=str(png_path), full_page=False, omit_background=False)
            screenshots.append(png_path)
            print(f"  [{i}/{total}] {png_path.name}", file=sys.stderr)

        browser.close()

    server.shutdown()
    return screenshots


def build_pptx(screenshots: list[Path], notes: list[dict], out_path: Path) -> None:
    """16:9 PPTX 생성, 슬라이드당 PNG 풀 블리드 + 노트 텍스트."""
    pres = Presentation()
    pres.slide_width = Emu(SLIDE_W_EMU)
    pres.slide_height = Emu(SLIDE_H_EMU)
    blank_layout = pres.slide_layouts[6]  # 빈 레이아웃

    for i, png in enumerate(screenshots):
        slide = pres.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(png),
            Emu(0),
            Emu(0),
            width=Emu(SLIDE_W_EMU),
            height=Emu(SLIDE_H_EMU),
        )
        if i < len(notes):
            n = notes[i]
            tf = slide.notes_slide.notes_text_frame
            header = " · ".join(x for x in (n.get("pnum"), n.get("pact"), n.get("title")) if x)
            body = n.get("script", "").strip()
            tf.text = f"{header}\n\n{body}" if header else body

    pres.save(str(out_path))


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("html", help="메인 덱 HTML 경로 (예: output/분기_OKR_리뷰/index.html)")
    ap.add_argument("-o", "--output", help="출력 PPTX 경로 (기본: HTML과 같은 이름의 .pptx)")
    ap.add_argument("--notes", help="스크립트 뷰어 HTML 경로 (기본: 같은 디렉토리의 스크립트_화면.html)")
    ap.add_argument("--keep-screenshots", action="store_true", help="중간 PNG 보존 (디버깅용)")
    ap.add_argument("--settle-ms", type=int, default=300, help="슬라이드 전환 후 정착 대기 ms (기본 300)")
    args = ap.parse_args()

    html = Path(args.html).resolve()
    if not html.exists():
        print(f"❌ HTML 파일 없음: {html}", file=sys.stderr)
        return 1

    out = Path(args.output).resolve() if args.output else html.with_suffix(".pptx")
    notes_html = Path(args.notes).resolve() if args.notes else html.parent / "스크립트_화면.html"

    print(f"📄 입력: {html}", file=sys.stderr)
    print(f"📝 노트: {notes_html} {'(있음)' if notes_html.exists() else '(없음 — 노트 비움)'}", file=sys.stderr)

    notes = parse_notes(notes_html)
    print(f"📝 추출된 노트: {len(notes)}개 (.page 슬라이드)", file=sys.stderr)

    with tempfile.TemporaryDirectory(prefix="pptx_screenshot_") as tmp:
        tmp_path = Path(tmp)
        print("📸 슬라이드 스크린샷 중...", file=sys.stderr)
        screenshots = screenshot_slides(html, tmp_path, settle_ms=args.settle_ms)

        if len(notes) and len(notes) != len(screenshots):
            print(
                f"⚠️  슬라이드 수({len(screenshots)}) ≠ 노트 수({len(notes)}). "
                f"앞에서부터 매칭, 뒷부분 미스매치는 노트 누락/덤핑.",
                file=sys.stderr,
            )

        print(f"📦 PPTX 생성: {out}", file=sys.stderr)
        out.parent.mkdir(parents=True, exist_ok=True)
        build_pptx(screenshots, notes, out)

        if args.keep_screenshots:
            kept = out.parent / f"{out.stem}_screenshots"
            if kept.exists():
                shutil.rmtree(kept)
            shutil.copytree(tmp_path, kept)
            print(f"🖼️  스크린샷 보존: {kept}", file=sys.stderr)

    print(f"✅ 완료: {out}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
